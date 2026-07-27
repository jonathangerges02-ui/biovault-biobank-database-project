"""Validate the final repository structure and binary deliverables."""

from __future__ import annotations

import json
import subprocess
import zipfile
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pypdf import PdfReader


ROOT = Path(__file__).resolve().parents[1]

REQUIRED = [
    "README.md",
    "report.docx",
    "report.pdf",
    "presentation.pptx",
    "sql/create_tables.sql",
    "sql/load_data.sql",
    "sql/queries.sql",
    "sql/views.sql",
    "sql/triggers_procedures.sql",
    "sql/test_constraints.sql",
    "diagrams/ERD.png",
    "diagrams/ERD.mmd",
    "src/app.py",
    "src/database.py",
    "src/repository.py",
    "presentation/BioVault_Presentation_Video.mp4",
    "presentation/voiceover_script.md",
    "presentation/live_demo_runbook.md",
    "presentation/defense_questions.md",
    "evidence/test_summary.md",
]


def require(condition: bool, message: str) -> None:
    if not condition:
        raise AssertionError(message)


def main() -> None:
    missing = [path for path in REQUIRED if not (ROOT / path).exists()]
    require(not missing, f"Missing required deliverables: {missing}")

    with zipfile.ZipFile(ROOT / "report.docx") as archive:
        require(archive.testzip() is None, "report.docx is a damaged ZIP package")
        require("word/document.xml" in archive.namelist(), "report.docx is not Word OOXML")

    with zipfile.ZipFile(ROOT / "presentation.pptx") as archive:
        require(archive.testzip() is None, "presentation.pptx is a damaged ZIP package")
        require(
            "ppt/presentation.xml" in archive.namelist(),
            "presentation.pptx is not PowerPoint OOXML",
        )

    pdf = PdfReader(ROOT / "report.pdf")
    require(len(pdf.pages) >= 10, "Report PDF should contain at least 10 pages")

    deck = Presentation(ROOT / "presentation.pptx")
    require(len(deck.slides) == 18, "Presentation must contain 18 slides")

    with Image.open(ROOT / "diagrams" / "ERD.png") as erd:
        require(erd.width >= 3000 and erd.height >= 1600, "ERD resolution is too low")

    video_path = ROOT / "presentation" / "BioVault_Presentation_Video.mp4"
    probe = subprocess.run(
        [
            "ffprobe",
            "-v",
            "error",
            "-show_entries",
            "format=duration",
            "-of",
            "json",
            str(video_path),
        ],
        check=True,
        capture_output=True,
        text=True,
    )
    duration = float(json.loads(probe.stdout)["format"]["duration"])
    require(900 <= duration <= 1200, f"Video duration {duration:.1f}s is not 15–20 minutes")

    ddl = (ROOT / "sql" / "create_tables.sql").read_text(encoding="utf-8")
    views = (ROOT / "sql" / "views.sql").read_text(encoding="utf-8")
    logic = (ROOT / "sql" / "triggers_procedures.sql").read_text(encoding="utf-8")
    require(ddl.count("CREATE TABLE ") == 15, "Expected 15 CREATE TABLE statements")
    require(ddl.count("CREATE INDEX ") == 13, "Expected 13 targeted CREATE INDEX statements")
    require(views.count("CREATE OR REPLACE VIEW ") == 3, "Expected three views")
    require("FOR UPDATE OF a" in logic, "Concurrency row lock is missing")
    require("RESEARCH_USE" in logic, "Consent enforcement is missing")

    print("Submission validation passed.")
    print(f"Required files: {len(REQUIRED)}")
    print(f"Report pages: {len(pdf.pages)}")
    print(f"Presentation slides: {len(deck.slides)}")
    print(f"Video duration: {duration:.1f} seconds")
    print("Schema objects: 15 tables, 13 targeted indexes, 3 views")


if __name__ == "__main__":
    main()
