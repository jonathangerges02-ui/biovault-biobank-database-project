"""Generate the PowerPoint, rendered slide frames, and voiceover script."""

from __future__ import annotations

import json
import textwrap
from pathlib import Path
from typing import Any

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
PRESENTATION_DIR = ROOT / "presentation"
SLIDES = json.loads((PRESENTATION_DIR / "slides.json").read_text(encoding="utf-8"))
RENDER_DIR = PRESENTATION_DIR / "rendered"

FONT_REGULAR = "C:/Windows/Fonts/arial.ttf"
FONT_BOLD = "C:/Windows/Fonts/arialbd.ttf"

NAV = "#123B4A"
ACCENT = "#2BAE9B"
BLUE = "#4285C5"
PURPLE = "#8C6BC1"
CORAL = "#D5656F"
INK = "#18323B"
MUTED = "#607982"
PAPER = "#F4F8F9"
WHITE = "#FFFFFF"
BORDER = "#C9DDE2"


def rgb(hex_color: str) -> RGBColor:
    value = hex_color.lstrip("#")
    return RGBColor.from_string(value)


def pil_font(size: int, bold: bool = False) -> ImageFont.FreeTypeFont:
    return ImageFont.truetype(FONT_BOLD if bold else FONT_REGULAR, size)


def set_shape_fill(shape: Any, color: str, transparency: int = 0) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(color)
    shape.fill.transparency = transparency
    shape.line.fill.background()


def add_text(
    slide: Any,
    text: str,
    left: float,
    top: float,
    width: float,
    height: float,
    *,
    size: int,
    color: str,
    bold: bool = False,
    font_name: str = "Arial",
    align: PP_ALIGN = PP_ALIGN.LEFT,
    valign: MSO_ANCHOR = MSO_ANCHOR.TOP,
) -> Any:
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    frame = box.text_frame
    frame.clear()
    frame.margin_left = frame.margin_right = Inches(0)
    frame.margin_top = frame.margin_bottom = Inches(0)
    frame.vertical_anchor = valign
    paragraph = frame.paragraphs[0]
    paragraph.text = text
    paragraph.alignment = align
    run = paragraph.runs[0]
    run.font.name = font_name
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = rgb(color)
    return box


def add_picture_contain(
    slide: Any,
    image_path: Path,
    left: float,
    top: float,
    width: float,
    height: float,
) -> None:
    with Image.open(image_path) as image:
        iw, ih = image.size
    image_ratio = iw / ih
    box_ratio = width / height
    if image_ratio > box_ratio:
        actual_width = width
        actual_height = width / image_ratio
        actual_left = left
        actual_top = top + (height - actual_height) / 2
    else:
        actual_height = height
        actual_width = height * image_ratio
        actual_left = left + (width - actual_width) / 2
        actual_top = top
    slide.shapes.add_picture(
        str(image_path),
        Inches(actual_left),
        Inches(actual_top),
        Inches(actual_width),
        Inches(actual_height),
    )


VISUAL_LABELS = {
    "problem": [("1", "Donor"), ("N", "Samples"), ("M:N", "Projects"), ("Δ", "Inventory")],
    "scope": [("IN", "Research"), ("OUT", "Direct ID"), ("WHO", "Lab teams"), ("DATA", "Synthetic")],
    "rules": [("✓", "Consent"), ("✓", "Membership"), ("✓", "Stock"), ("✓", "Audit")],
    "cardinality": [("1", "parent"), ("0..*", "children"), ("1..*", "required"), ("M:N", "bridge")],
    "schema": [("PK", "identity"), ("UK", "codes"), ("FK", "lineage"), ("CK", "domains")],
    "normalization": [("1NF", "atomic"), ("2NF", "whole key"), ("3NF", "no transitives"), ("JSONB", "history")],
    "implementation": [("01", "DDL"), ("02", "logic"), ("03", "data"), ("04", "views")],
    "constraints": [("CHECK", "valid states"), ("UNIQUE", "identity"), ("FK", "lineage"), ("INDEX", "access")],
    "queries": [("JOIN", "lineage"), ("Σ", "aggregate"), ("WITH", "recursive"), ("CRUD", "safe")],
    "data": [("12", "donors"), ("20", "samples"), ("15", "tests"), ("12", "usage")],
    "demo": [("1", "setup"), ("2", "tests"), ("3", "queries"), ("4", "UI")],
    "future": [("RBAC", "security"), ("BAR", "barcodes"), ("IoT", "sensors"), ("CoC", "custody")],
    "conclusion": [("ERD", "designed"), ("3NF", "mapped"), ("SQL", "verified"), ("UI", "working")],
}


def add_four_cards_pptx(slide: Any, visual: str) -> None:
    labels = VISUAL_LABELS.get(visual, [("01", "Design"), ("02", "Build"), ("03", "Test"), ("04", "Defend")])
    colors = (ACCENT, BLUE, PURPLE, CORAL)
    positions = ((7.20, 2.05), (9.88, 2.05), (7.20, 4.22), (9.88, 4.22))
    for (big, small), color, (left, top) in zip(labels, colors, positions):
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(left),
            Inches(top),
            Inches(2.35),
            Inches(1.72),
        )
        set_shape_fill(card, WHITE)
        card.line.color.rgb = rgb(BORDER)
        card.line.width = Pt(1.4)
        accent = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(left),
            Inches(top),
            Inches(0.10),
            Inches(1.72),
        )
        set_shape_fill(accent, color)
        add_text(slide, big, left + 0.25, top + 0.25, 1.9, 0.65, size=26, color=color, bold=True)
        add_text(slide, small, left + 0.25, top + 1.02, 1.9, 0.42, size=15, color=MUTED, bold=True)


def add_trigger_flow_pptx(slide: Any) -> None:
    items = [
        ("1", "Lock", "FOR UPDATE", ACCENT),
        ("2", "Authorize", "project + member", BLUE),
        ("3", "Consent", "active on date", PURPLE),
        ("4", "Deduct", "atomic stock", CORAL),
    ]
    top = 1.94
    for index, (number, title, detail, color) in enumerate(items):
        left = 7.10 + (index % 2) * 2.85
        row_top = top + (index // 2) * 2.16
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(left),
            Inches(row_top),
            Inches(2.50),
            Inches(1.75),
        )
        set_shape_fill(shape, WHITE)
        shape.line.color.rgb = rgb(BORDER)
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(left + 0.18),
            Inches(row_top + 0.20),
            Inches(0.55),
            Inches(0.55),
        )
        set_shape_fill(circle, color)
        add_text(slide, number, left + 0.18, row_top + 0.30, 0.55, 0.25, size=14, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, title, left + 0.87, row_top + 0.20, 1.35, 0.45, size=19, color=INK, bold=True)
        add_text(slide, detail, left + 0.30, row_top + 1.02, 1.95, 0.42, size=14, color=MUTED)


def add_visual_pptx(slide: Any, visual: str) -> None:
    image_map = {
        "erd": ROOT / "diagrams" / "ERD.png",
        "tests": ROOT / "evidence" / "database_tests.png",
        "ui": ROOT / "evidence" / "ui_dashboard.png",
    }
    if visual in image_map:
        frame = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(6.82),
            Inches(1.60),
            Inches(5.92),
            Inches(4.98),
        )
        set_shape_fill(frame, WHITE)
        frame.line.color.rgb = rgb(BORDER)
        add_picture_contain(slide, image_map[visual], 6.98, 1.76, 5.60, 4.66)
    elif visual == "trigger_flow":
        add_trigger_flow_pptx(slide)
    else:
        add_four_cards_pptx(slide, visual)


def add_slide_chrome(slide: Any, index: int, title: str, kicker: str) -> None:
    background = slide.background
    background.fill.solid()
    background.fill.fore_color.rgb = rgb(PAPER)
    band = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(0.16)
    )
    set_shape_fill(band, ACCENT)
    add_text(slide, kicker.upper(), 0.65, 0.43, 11.7, 0.3, size=10, color=ACCENT, bold=True)
    add_text(slide, title, 0.65, 0.82, 11.9, 0.64, size=28, color=NAV, bold=True)
    add_text(slide, "BioVault", 0.65, 7.12, 1.4, 0.20, size=9, color=MUTED, bold=True)
    add_text(slide, f"{index:02d}", 12.15, 7.12, 0.55, 0.20, size=9, color=MUTED, bold=True, align=PP_ALIGN.RIGHT)


def add_bullets_pptx(slide: Any, bullets: list[str]) -> None:
    box = slide.shapes.add_textbox(Inches(0.72), Inches(1.82), Inches(5.62), Inches(4.75))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Inches(0)
    frame.margin_right = Inches(0.06)
    for index, bullet in enumerate(bullets):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.text = bullet
        paragraph.level = 0
        paragraph.space_after = Pt(20)
        paragraph.font.name = "Arial"
        paragraph.font.size = Pt(21)
        paragraph.font.color.rgb = rgb(INK)
        paragraph.font.bold = index == 0
        paragraph.text = f"•  {bullet}"


def create_powerpoint() -> None:
    presentation = Presentation()
    presentation.slide_width = Inches(13.333)
    presentation.slide_height = Inches(7.5)
    blank = presentation.slide_layouts[6]

    for index, spec in enumerate(SLIDES, start=1):
        slide = presentation.slides.add_slide(blank)
        if index == 1:
            background = slide.background
            background.fill.solid()
            background.fill.fore_color.rgb = rgb(NAV)
            for offset, color in ((0.0, ACCENT), (0.32, BLUE), (0.64, PURPLE)):
                for row in range(5):
                    circle = slide.shapes.add_shape(
                        MSO_SHAPE.OVAL,
                        Inches(9.65 + offset + (row % 2) * 0.48),
                        Inches(0.62 + row * 1.25),
                        Inches(0.30),
                        Inches(0.30),
                    )
                    set_shape_fill(circle, color)
            add_text(slide, "BIOVAULT", 0.82, 0.72, 3.2, 0.35, size=12, color=ACCENT, bold=True)
            add_text(slide, spec["title"], 0.82, 1.62, 7.9, 0.82, size=44, color=WHITE, bold=True)
            add_text(slide, spec["subtitle"], 0.82, 2.58, 8.5, 0.58, size=25, color="#CDE3E8")
            line = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, Inches(0.82), Inches(3.55), Inches(1.15), Inches(0.08)
            )
            set_shape_fill(line, ACCENT)
            add_text(slide, spec["kicker"], 0.82, 3.86, 8.1, 0.36, size=13, color="#A9CAD1", bold=True)
            add_text(slide, spec["bullets"][0], 0.82, 5.35, 4.0, 0.34, size=15, color=WHITE, bold=True)
            add_text(slide, spec["bullets"][1], 0.82, 5.78, 4.0, 0.30, size=13, color="#A9CAD1")
            add_text(slide, "01", 12.05, 7.12, 0.6, 0.2, size=9, color="#8AB2BB", bold=True, align=PP_ALIGN.RIGHT)
        else:
            add_slide_chrome(slide, index, spec["title"], spec.get("kicker", ""))
            add_bullets_pptx(slide, spec["bullets"])
            add_visual_pptx(slide, spec.get("visual", ""))

    presentation.save(ROOT / "presentation.pptx")


def contain_rect(iw: int, ih: int, x: int, y: int, w: int, h: int) -> tuple[int, int, int, int]:
    ratio = min(w / iw, h / ih)
    nw, nh = int(iw * ratio), int(ih * ratio)
    return x + (w - nw) // 2, y + (h - nh) // 2, nw, nh


def render_visual_pil(draw: ImageDraw.ImageDraw, canvas: Image.Image, visual: str) -> None:
    image_map = {
        "erd": ROOT / "diagrams" / "ERD.png",
        "tests": ROOT / "evidence" / "database_tests.png",
        "ui": ROOT / "evidence" / "ui_dashboard.png",
    }
    if visual in image_map:
        draw.rounded_rectangle((665, 158, 1225, 640), radius=18, fill=WHITE, outline=BORDER, width=2)
        with Image.open(image_map[visual]).convert("RGB") as source:
            rect = contain_rect(*source.size, 682, 174, 526, 450)
            source.thumbnail((rect[2], rect[3]), Image.Resampling.LANCZOS)
            canvas.paste(source, (rect[0], rect[1]))
        return

    labels = VISUAL_LABELS.get(visual, [("01", "Design"), ("02", "Build"), ("03", "Test"), ("04", "Defend")])
    colors = (ACCENT, BLUE, PURPLE, CORAL)
    positions = ((690, 205), (950, 205), (690, 420), (950, 420))
    for (big, small), color, (x, y) in zip(labels, colors, positions):
        draw.rounded_rectangle((x, y, x + 230, y + 165), radius=18, fill=WHITE, outline=BORDER, width=2)
        draw.rectangle((x, y, x + 9, y + 165), fill=color)
        draw.text((x + 28, y + 26), big, font=pil_font(29, True), fill=color)
        draw.text((x + 28, y + 107), small, font=pil_font(18, True), fill=MUTED)


def render_frames() -> None:
    RENDER_DIR.mkdir(parents=True, exist_ok=True)
    for index, spec in enumerate(SLIDES, start=1):
        canvas = Image.new("RGB", (1280, 720), NAV if index == 1 else PAPER)
        draw = ImageDraw.Draw(canvas)
        if index == 1:
            draw.text((80, 68), "BIOVAULT", font=pil_font(16, True), fill=ACCENT)
            draw.text((80, 160), spec["title"], font=pil_font(58, True), fill=WHITE)
            draw.text((80, 246), spec["subtitle"], font=pil_font(30), fill="#CDE3E8")
            draw.rectangle((80, 352, 195, 360), fill=ACCENT)
            draw.text((80, 390), spec["kicker"], font=pil_font(18, True), fill="#A9CAD1")
            draw.text((80, 548), spec["bullets"][0], font=pil_font(21, True), fill=WHITE)
            draw.text((80, 588), spec["bullets"][1], font=pil_font(18), fill="#A9CAD1")
            for row in range(5):
                for col, color in enumerate((ACCENT, BLUE, PURPLE)):
                    x = 1000 + col * 45 + (row % 2) * 22
                    y = 95 + row * 112
                    draw.ellipse((x, y, x + 26, y + 26), fill=color)
            draw.text((1180, 682), "01", font=pil_font(13, True), fill="#8AB2BB")
        else:
            draw.rectangle((0, 0, 1280, 15), fill=ACCENT)
            draw.text((62, 43), spec.get("kicker", "").upper(), font=pil_font(14, True), fill=ACCENT)
            draw.text((62, 83), spec["title"], font=pil_font(38, True), fill=NAV)
            y = 184
            for bullet in spec["bullets"]:
                wrapped = textwrap.wrap(bullet, width=39)
                draw.ellipse((70, y + 8, 82, y + 20), fill=ACCENT)
                draw.multiline_text((100, y), "\n".join(wrapped), font=pil_font(23), fill=INK, spacing=7)
                y += 56 + (len(wrapped) - 1) * 31
            render_visual_pil(draw, canvas, spec.get("visual", ""))
            draw.text((62, 680), "BioVault", font=pil_font(13, True), fill=MUTED)
            draw.text((1185, 680), f"{index:02d}", font=pil_font(13, True), fill=MUTED)
        canvas.save(RENDER_DIR / f"slide_{index:02d}.png", quality=95)


def write_voiceover_script() -> None:
    total_seconds = sum(int(spec["duration_seconds"]) for spec in SLIDES)
    lines = [
        "# BioVault voiceover script",
        "",
        f"Planned duration: **{total_seconds // 60}:{total_seconds % 60:02d}** "
        "(18 slides). Speak naturally and replace the student placeholders.",
        "",
        "> The generated rehearsal video uses a synthetic system voice. For academic "
        "submission, record this script in the submitting student's own voice if the "
        "instructor requires personal narration.",
        "",
    ]
    elapsed = 0
    for index, spec in enumerate(SLIDES, start=1):
        start = f"{elapsed // 60}:{elapsed % 60:02d}"
        elapsed += int(spec["duration_seconds"])
        end = f"{elapsed // 60}:{elapsed % 60:02d}"
        lines.extend(
            [
                f"## Slide {index}: {spec['title']} ({start}–{end})",
                "",
                spec["narration"],
                "",
            ]
        )
    (PRESENTATION_DIR / "voiceover_script.md").write_text(
        "\n".join(lines), encoding="utf-8"
    )


if __name__ == "__main__":
    create_powerpoint()
    render_frames()
    write_voiceover_script()
    print("Generated presentation.pptx, 18 rendered slides, and voiceover script.")
